import hashlib
import csv
import io
import re
from pathlib import Path
from typing import List

from pypdf import PdfReader

from rag_book_agent.models import Page, ParsedDocument


class UnsupportedDocumentError(ValueError):
    pass


class DocumentLoader:
    supported_extensions = {
        ".pdf", ".md", ".markdown", ".txt", ".docx", ".xlsx", ".xls", ".pptx", ".csv"
    }

    def __init__(self, settings=None):
        self.settings = settings

    def load(self, path: Path) -> ParsedDocument:
        path = path.expanduser().resolve()
        if not path.is_file():
            raise FileNotFoundError("Document does not exist: %s" % path)

        extension = path.suffix.lower()
        if extension not in self.supported_extensions:
            raise UnsupportedDocumentError("Unsupported document type: %s" % extension)

        content_hash = self._hash_file(path)
        if extension == ".pdf":
            pages = self._load_pdf(path)
            file_type = "pdf"
        elif extension in {".md", ".markdown", ".txt", ".csv"}:
            pages = self._load_text(path)
            if extension == ".csv":
                pages = self._load_csv(path)
            file_type = "markdown" if extension in {".md", ".markdown"} else extension[1:]
        elif extension == ".docx":
            pages, file_type = self._load_docx(path), "docx"
        elif extension == ".xlsx":
            pages, file_type = self._load_xlsx(path), "excel"
        elif extension == ".xls":
            pages, file_type = self._load_xls(path), "excel"
        elif extension == ".pptx":
            pages, file_type = self._load_pptx(path), "pptx"

        title = path.stem.replace("_", " ").strip()
        return ParsedDocument(
            path=str(path),
            title=title,
            file_type=file_type,
            content_hash=content_hash,
            pages=pages,
        )

    def find_documents(self, path: Path) -> List[Path]:
        path = path.expanduser().resolve()
        if path.is_file():
            return [path]
        if not path.is_dir():
            raise FileNotFoundError("Path does not exist: %s" % path)

        documents = []
        for item in path.rglob("*"):
            if item.is_file() and item.suffix.lower() in self.supported_extensions:
                documents.append(item)
        return sorted(documents)

    def _load_pdf(self, path: Path) -> List[Page]:
        # Keep page boundaries and coordinates through PyMuPDF first.
        try:
            pages = self._load_pdf_pymupdf(path)
            if all(len(page.text) >= self._min_text_chars() for page in pages):
                return pages
        except Exception:
            pass
        # Scanned pages need a real OCR/layout engine. It is lazy-loaded.
        if self._ocr_enabled():
            try:
                pages = self._load_pdf_pymupdf_with_ocr(path)
                if any(page.text for page in pages):
                    return pages
            except Exception as exc:
                # Keep a diagnostic page so callers can report the real failure.
                return [Page(number=1, text="", parser="ocr_error", ocr_error=str(exc))]
        # MarkItDown and pypdf remain compatibility fallbacks.
        try:
            pages = self._load_pdf_markitdown(path)
            if any(page.text for page in pages):
                return pages
        except Exception:
            pass
        return self._load_pdf_pypdf(path)

    def _load_pdf_pymupdf_with_ocr(self, path: Path) -> List[Page]:
        import fitz
        from rag_book_agent.ingest.pdf_ocr import PaddlePdfReader

        document = fitz.open(str(path))
        reader = PaddlePdfReader(
            dpi=self._setting("pdf_render_dpi", 220),
            formula=self._setting("pdf_formula_enabled", True),
            tables=self._setting("pdf_table_enabled", True),
        )
        pages = []
        try:
            for number, pdf_page in enumerate(document, start=1):
                native = self._normalize("\n".join(block[4] for block in pdf_page.get_text("blocks", sort=True) if len(block) >= 5 and block[4]))
                if len(native) >= self._min_text_chars():
                    pages.append(Page(number=number, text=native, parser="pymupdf", image_count=len(pdf_page.get_images())))
                    continue
                text, confidence, table_count = reader.read_page(pdf_page)
                pages.append(Page(number=number, text=self._normalize(text), parser="paddleocr", ocr_used=True, ocr_confidence=confidence, image_count=len(pdf_page.get_images()), table_count=table_count))
        finally:
            document.close()
        return pages

    def _setting(self, name: str, default):
        return getattr(self.settings, name, default) if self.settings is not None else default

    def _min_text_chars(self) -> int:
        return int(self._setting("pdf_min_text_chars", 20))

    def _ocr_enabled(self) -> bool:
        return bool(self._setting("pdf_ocr_enabled", False))

    def _load_pdf_markitdown(self, path: Path) -> List[Page]:
        from markitdown import MarkItDown

        result = MarkItDown().convert(str(path))
        text = self._normalize(getattr(result, "text_content", "") or "")
        raw_pages = [part for part in text.split("\f") if part.strip()]
        if not raw_pages:
            raw_pages = [text]
        return [Page(number=index, text=part) for index, part in enumerate(raw_pages, start=1)]

    def _load_pdf_pymupdf(self, path: Path) -> List[Page]:
        import fitz

        document = fitz.open(str(path))
        pages = []
        try:
            for number, pdf_page in enumerate(document, start=1):
                blocks = pdf_page.get_text("blocks", sort=True)
                text_parts = []
                for block in blocks:
                    if len(block) >= 5 and block[4]:
                        text_parts.append(block[4])
                text = self._normalize("\n".join(text_parts))
                pages.append(Page(number=number, text=text))
        finally:
            document.close()
        return pages

    def _load_pdf_pypdf(self, path: Path) -> List[Page]:
        reader = PdfReader(str(path), strict=False)
        pages = []
        for number, pdf_page in enumerate(reader.pages, start=1):
            try:
                text = pdf_page.extract_text() or ""
            except (OSError, ValueError, RuntimeError):
                text = ""
            pages.append(Page(number=number, text=self._normalize(text)))
        return pages

    def _load_text(self, path: Path) -> List[Page]:
        try:
            text = path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            text = path.read_text(encoding="utf-8-sig")
        return [Page(number=1, text=self._normalize(text))]

    def _load_csv(self, path: Path) -> List[Page]:
        text = path.read_text(encoding="utf-8-sig")
        rows = list(csv.reader(io.StringIO(text)))
        if not rows:
            return [Page(number=1, text="")]
        width = max(len(row) for row in rows)
        header = rows[0] + [""] * (width - len(rows[0]))
        lines = ["| " + " | ".join(cell.replace("|", "\\|") for cell in header) + " |"]
        lines.append("| " + " | ".join("---" for _ in range(width)) + " |")
        for row in rows[1:]:
            values = row + [""] * (width - len(row))
            lines.append("| " + " | ".join(cell.replace("|", "\\|") for cell in values) + " |")
        return [Page(number=1, text="\n".join(lines))]

    def _load_docx(self, path: Path) -> List[Page]:
        from docx import Document

        document = Document(str(path))
        lines = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables:
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            if rows:
                lines.append("\n".join("| " + " | ".join(row) + " |" for row in rows))
        return [Page(number=1, text=self._normalize("\n\n".join(lines)))]

    def _load_xlsx(self, path: Path) -> List[Page]:
        from openpyxl import load_workbook

        workbook = load_workbook(str(path), read_only=True, data_only=True)
        pages = []
        try:
            for number, sheet in enumerate(workbook.worksheets, start=1):
                rows = list(sheet.iter_rows(values_only=True))
                lines = ["## " + sheet.title]
                for row in rows:
                    values = ["" if value is None else str(value) for value in row]
                    lines.append("| " + " | ".join(values) + " |")
                pages.append(Page(number=number, text=self._normalize("\n".join(lines))))
        finally:
            workbook.close()
        return pages or [Page(number=1, text="")]

    def _load_xls(self, path: Path) -> List[Page]:
        import xlrd

        workbook = xlrd.open_workbook(str(path), on_demand=True)
        pages = []
        try:
            for number, sheet in enumerate(workbook.sheets(), start=1):
                lines = ["## " + sheet.name]
                for row_index in range(sheet.nrows):
                    values = [str(value) for value in sheet.row_values(row_index)]
                    lines.append("| " + " | ".join(values) + " |")
                pages.append(Page(number=number, text=self._normalize("\n".join(lines))))
        finally:
            workbook.release_resources()
        return pages or [Page(number=1, text="")]

    def _load_pptx(self, path: Path) -> List[Page]:
        from pptx import Presentation

        presentation = Presentation(str(path))
        pages = []
        for number, slide in enumerate(presentation.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            pages.append(Page(number=number, text=self._normalize("\n\n".join(texts))))
        return pages or [Page(number=1, text="")]

    @staticmethod
    def _normalize(text: str) -> str:
        text = text.replace("\r\n", "\n").replace("\r", "\n")
        # Join words split at a line ending while keeping paragraph boundaries.
        text = re.sub(r"(?<=\w)-\n(?=\w)", "", text)
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()

    @staticmethod
    def _hash_file(path: Path) -> str:
        digest = hashlib.sha256()
        with path.open("rb") as handle:
            while True:
                block = handle.read(1024 * 1024)
                if not block:
                    break
                digest.update(block)
        return digest.hexdigest()
